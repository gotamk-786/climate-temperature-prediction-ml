from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
SLIDES_DIR = PROJECT_ROOT / "slides"
FIGURES_DIR = PROJECT_ROOT / "outputs" / "figures"
TABLES_DIR = PROJECT_ROOT / "outputs" / "tables"
OUTPUT_PATH = SLIDES_DIR / "climate_change_prediction_presentation.pptx"

BG = RGBColor(248, 250, 247)
INK = RGBColor(31, 41, 35)
MUTED = RGBColor(91, 104, 94)
GREEN = RGBColor(22, 104, 76)
GOLD = RGBColor(184, 130, 47)
LINE = RGBColor(210, 218, 207)
WHITE = RGBColor(255, 255, 255)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_band(slide, color: RGBColor = GREEN) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.0), Inches(0.65))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.98), Inches(11.5), Inches(0.34))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Aptos"
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def add_footer(slide, number: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.05), Inches(0.9), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = f"{number:02d}"
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_body(slide, bullets: list[str], x: float, y: float, w: float, h: float, font_size: int = 17) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_metric(slide, label: str, value: str, x: float, y: float, w: float = 2.3) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.name = "Aptos Display"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = "Aptos"
    p2.font.size = Pt(10)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER


def add_callout(slide, title: str, text: str, x: float, y: float, w: float, h: float = 1.05, color: RGBColor = GREEN) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Aptos Display"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.name = "Aptos"
    p2.font.size = Pt(10)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER


def add_pipeline(slide, steps: list[str], x: float, y: float, w: float) -> None:
    step_w = w / len(steps)
    for i, step in enumerate(steps):
        sx = x + i * step_w
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx), Inches(y), Inches(step_w - 0.18), Inches(0.82))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LINE
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = step
        p.font.name = "Aptos"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = INK
        p.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(sx + step_w - 0.28), Inches(y + 0.23), Inches(0.28), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()


def add_picture(slide, filename: str, x: float, y: float, w: float, h: float | None = None) -> None:
    path = FIGURES_DIR / filename
    if h is None:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    else:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_table(slide, df: pd.DataFrame, x: float, y: float, w: float, h: float, font_size: int = 10) -> None:
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            p.font.name = "Aptos"
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    for r in range(df.shape[0]):
        for c in range(cols):
            cell = table.cell(r + 1, c)
            value = df.iloc[r, c]
            if isinstance(value, float):
                text = f"{value:.3f}"
            else:
                text = str(value)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.color.rgb = INK
                p.alignment = PP_ALIGN.CENTER
    for c in range(cols):
        table.columns[c].width = int(Inches(w / cols))


def new_slide(prs: Presentation, number: int, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_band(slide)
    add_title(slide, title, subtitle)
    add_footer(slide, number)
    return slide


def main() -> None:
    model_results = pd.read_csv(TABLES_DIR / "model_comparison_results.csv")
    cv_summary = pd.read_csv(TABLES_DIR / "cross_validation_summary.csv")
    error_by_decade = pd.read_csv(TABLES_DIR / "error_by_decade.csv")
    split_summary = pd.read_csv(TABLES_DIR / "split_strategy_summary.csv")
    tuning = pd.read_csv(TABLES_DIR / "tuning_before_after_comparison.csv")
    features = pd.read_csv(TABLES_DIR / "best_model_feature_importance.csv")
    warming = pd.read_csv(TABLES_DIR / "south_asia_warming_rates.csv")

    chrono = model_results[model_results["Split Strategy"] == "chronological"].copy()
    best = chrono.sort_values("RMSE").iloc[0]
    xgb = chrono[chrono["Model"] == "XGBoost"].iloc[0]
    lr = chrono[chrono["Model"] == "Linear Regression"].iloc[0]
    random_avg = split_summary[split_summary["Split Strategy"] == "random"].iloc[0]
    chrono_avg = split_summary[split_summary["Split Strategy"] == "chronological"].iloc[0]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = new_slide(prs, 1, "Climate Change Analysis and Temperature Prediction", "Predicting city-level temperature trends using historical climate records and machine learning")
    add_body(
        slide,
        [
            "Student: Gotam Kumar (23k-0860)",
            "Course: Data Science",
            "Instructor: Dr. Rabia Tabassum",
            "Institution: FAST-NUCES",
            "Dataset: Berkeley Earth GlobalLandTemperaturesByCity.csv",
        ],
        0.75,
        1.7,
        5.8,
        3.2,
        18,
    )
    add_picture(slide, "global_temperature_trend.png", 6.6, 1.55, 5.85, 2.35)
    add_metric(slide, "Best RMSE", f"{best['RMSE']:.3f}", 6.7, 4.35)
    add_metric(slide, "Best R2", f"{best['R2 Score']:.3f}", 9.2, 4.35)
    add_callout(slide, "Final Recommendation", "Random Forest for chronological prediction", 6.7, 5.65, 4.8, 0.78, GOLD)

    slide = new_slide(prs, 2, "Problem and Research Goal")
    add_body(
        slide,
        [
            "Climate change affects agriculture, health, water security, infrastructure, and regional planning.",
            "The project predicts monthly city-level AverageTemperature from historical climate records.",
            "This is a regression task because the target is a continuous temperature value, not a class label.",
            "Research question: which feature engineering and model strategy gives the best balance of accuracy, generalization, and climate relevance?",
        ],
        0.75,
        1.55,
        6.0,
        4.8,
        18,
    )
    add_picture(slide, "anomaly_trend_by_decade.png", 7.0, 1.55, 5.5, 2.15)
    add_picture(slide, "global_vs_south_asia_trend.png", 7.0, 4.05, 5.5, 2.15)

    slide = new_slide(prs, 3, "End-to-End Workflow", "How the code moves from raw climate data to final evaluation")
    add_pipeline(slide, ["Raw Data", "Cleaning", "Features", "EDA", "Models", "Evaluation", "Paper/PPT"], 0.75, 1.5, 11.7)
    add_body(
        slide,
        [
            "Raw CSV is loaded from data/raw/GlobalLandTemperaturesByCity.csv.",
            "Cleaning and feature engineering are handled in src/preprocessing.py.",
            "EDA figures are generated into outputs/figures.",
            "Modeling and tuning produce CSV result tables in outputs/tables.",
            "The research paper and this presentation are generated from the saved artifacts.",
        ],
        0.9,
        2.75,
        5.6,
        3.6,
        17,
    )
    add_picture(slide, "model_comparison_bar_chart.png", 6.9, 2.65, 5.4, 2.35)
    add_picture(slide, "feature_importance_best_model.png", 6.9, 5.15, 5.4, 1.55)

    slide = new_slide(prs, 4, "Dataset")
    add_body(
        slide,
        [
            "Source: Berkeley Earth Climate Change Dataset on Kaggle.",
            "File used: GlobalLandTemperaturesByCity.csv.",
            "Raw scale: approximately 8.6 million monthly city-level records.",
            "Columns: dt, AverageTemperature, AverageTemperatureUncertainty, City, Country, Latitude, Longitude.",
            "Working analysis window: 1900 onward after cleaning.",
        ],
        0.75,
        1.45,
        5.8,
        4.7,
        17,
    )
    add_picture(slide, "missing_values_chart.png", 6.8, 1.45, 5.4, 2.35)
    add_picture(slide, "temperature_distribution.png", 6.8, 4.05, 5.4, 2.35)

    slide = new_slide(prs, 5, "Preprocessing Pipeline", "Implemented in src/preprocessing.py and notebooks/01_data_cleaning.ipynb")
    add_body(
        slide,
        [
            "Parsed dt as datetime and converted temperature columns to numeric.",
            "Dropped rows with missing dt or AverageTemperature and removed duplicate records.",
            "Converted latitude/longitude strings into signed float values.",
            "Restricted the working dataset to 1900 onward.",
            "Created cleaned data, South Asia subset, modeling sample, and chronological train/test split.",
        ],
        0.75,
        1.55,
        6.1,
        4.9,
        17,
    )
    add_metric(slide, "Cleaned Rows", "4.79M", 7.1, 1.7)
    add_metric(slide, "Modeling Rows", "75K", 9.6, 1.7)
    add_metric(slide, "Train Rows", "60K", 7.1, 3.15)
    add_metric(slide, "Test Rows", "15K", 9.6, 3.15)

    slide = new_slide(prs, 6, "Feature Engineering")
    add_body(
        slide,
        [
            "Temporal: Year, Month, Quarter, Season, Decade, YearsSince1900.",
            "Geographic: LatitudeValue, LongitudeValue, Hemisphere, Country, RegionTag.",
            "Historical signals: Lag1Temperature, Lag12Temperature, Rolling12MeanTemperature.",
            "Climate interpretation: HistoricalCityMonthMean, TemperatureAnomaly, ClimateRiskIndex.",
            "South Asia region tag supports local climate interpretation and viva discussion.",
        ],
        0.75,
        1.55,
        6.0,
        4.8,
        17,
    )
    feat_table = features.head(6).copy()
    feat_table["Importance"] = feat_table["Importance"].round(4)
    add_table(slide, feat_table, 7.0, 1.55, 5.5, 3.3, 8)
    add_picture(slide, "feature_importance_best_model.png", 7.0, 5.05, 5.5, 1.55)

    slide = new_slide(prs, 7, "EDA Findings")
    add_body(
        slide,
        [
            "Global trend shows a clear warming direction over the modern climate record.",
            "Seasonal boxplot confirms strong month/season structure in the target.",
            "Correlation heatmap supports the importance of lagged and historical average features.",
            "Outlier and country plots show wide geographic temperature variation.",
        ],
        0.75,
        1.5,
        5.0,
        4.8,
        17,
    )
    add_picture(slide, "seasonal_boxplot.png", 6.0, 1.45, 3.1, 2.1)
    add_picture(slide, "correlation_heatmap.png", 9.35, 1.45, 3.0, 2.5)
    add_picture(slide, "country_comparison.png", 6.0, 4.15, 6.35, 2.15)

    slide = new_slide(prs, 8, "Implemented Models", "Implemented in src/train.py and scripts/run_modeling.py")
    add_body(
        slide,
        [
            "Linear Regression: interpretable baseline.",
            "Decision Tree Regressor: non-linear rule-based model.",
            "Random Forest Regressor: robust bagging ensemble.",
            "Gradient Boosting Regressor: sequential boosting ensemble.",
            "XGBoost Regressor: optimized tree boosting model.",
            "Numerical features were imputed and scaled; categorical features were imputed and one-hot encoded.",
        ],
        0.75,
        1.45,
        6.3,
        5.3,
        17,
    )
    add_picture(slide, "model_comparison_bar_chart.png", 7.15, 1.55, 5.1, 2.4)
    add_metric(slide, "Models Compared", "5", 7.35, 4.55)
    add_metric(slide, "Validation Splits", "2", 9.85, 4.55)

    slide = new_slide(prs, 9, "Chronological Model Results")
    result_table = chrono[["Model", "MAE", "RMSE", "MAPE", "R2 Score"]].copy()
    result_table = result_table[result_table["Model"].isin(["Linear Regression", "Decision Tree", "Random Forest", "Gradient Boosting", "XGBoost"])]
    result_table = result_table.sort_values("RMSE")
    add_table(slide, result_table, 0.75, 1.35, 7.2, 3.6, 10)
    add_body(
        slide,
        [
            f"Best chronological model: {best['Model']} with RMSE {best['RMSE']:.3f} and R2 {best['R2 Score']:.3f}.",
            f"XGBoost is practically close on chronological split with RMSE {xgb['RMSE']:.3f}.",
            f"Linear Regression remains a useful baseline but has higher RMSE {lr['RMSE']:.3f}.",
            "MAPE is not the primary selection metric because near-zero temperatures inflate percentage errors.",
        ],
        8.35,
        1.45,
        4.2,
        4.8,
        16,
    )
    add_callout(slide, "MAPE Note", "High MAPE is caused by near-zero temperatures; RMSE and R2 are primary.", 8.35, 5.55, 4.1, 0.78, GOLD)

    slide = new_slide(prs, 10, "Cross-Validation and Generalization")
    cv_table = cv_summary[["Model", "CV Mean RMSE", "CV Mean R2"]].head(5).copy()
    add_table(slide, cv_table, 0.75, 1.4, 5.8, 3.35, 10)
    add_body(
        slide,
        [
            f"XGBoost has best average CV RMSE: {cv_summary.iloc[0]['CV Mean RMSE']:.4f}.",
            "Random Forest leads on chronological holdout, showing that ranking depends on validation design.",
            "Chronological validation is more realistic for climate prediction because it tests later years using earlier observations.",
        ],
        7.0,
        1.6,
        5.35,
        3.8,
        17,
    )
    add_metric(slide, "Random Split Avg RMSE", f"{random_avg['RMSE']:.4f}", 7.05, 5.25)
    add_metric(slide, "Chronological Avg RMSE", f"{chrono_avg['RMSE']:.4f}", 9.7, 5.25)

    slide = new_slide(prs, 11, "Comparative Analysis")
    add_callout(slide, "Random Forest", f"Best chronological RMSE: {best['RMSE']:.3f}", 0.85, 1.5, 3.0, 1.0, GREEN)
    add_callout(slide, "XGBoost", f"Best CV RMSE: {cv_summary.iloc[0]['CV Mean RMSE']:.4f}", 4.05, 1.5, 3.0, 1.0, GOLD)
    add_callout(slide, "Linear Regression", f"Weakest baseline RMSE: {lr['RMSE']:.3f}", 7.25, 1.5, 3.0, 1.0, RGBColor(99, 112, 102))
    add_body(
        slide,
        [
            f"Random Forest and XGBoost are practically close on chronological holdout: {best['RMSE']:.3f} vs {xgb['RMSE']:.3f}.",
            "Ensemble models outperform simple linear and single-tree baselines.",
            "XGBoost is stronger under cross-validation, while Random Forest is stronger under chronological testing.",
            "Decision Tree is easier to explain but weaker in predictive accuracy.",
            "Final recommendation: Random Forest for future-period prediction; XGBoost as the CV-focused alternative.",
        ],
        0.95,
        3.05,
        6.1,
        3.6,
        16,
    )
    add_picture(slide, "model_comparison_bar_chart.png", 7.35, 3.0, 5.0, 2.45)
    add_picture(slide, "actual_vs_predicted_best_model.png", 10.2, 1.3, 2.15, 2.15)

    slide = new_slide(prs, 12, "Best Model Diagnostics")
    add_picture(slide, "actual_vs_predicted_best_model.png", 0.8, 1.35, 4.6, 4.6)
    add_picture(slide, "residual_plot.png", 5.8, 1.4, 3.2, 2.1)
    add_picture(slide, "feature_importance_best_model.png", 5.8, 4.0, 6.2, 2.25)
    add_body(
        slide,
        [
            "Actual vs predicted plot shows close alignment around the ideal diagonal.",
            "Residuals are concentrated near zero.",
            "HistoricalCityMonthMean is the dominant feature, followed by lag and rolling features.",
        ],
        9.25,
        1.55,
        3.0,
        2.1,
        14,
    )

    slide = new_slide(prs, 13, "Hyperparameter Tuning")
    tuning_table = tuning[["Baseline Model", "Baseline RMSE", "Tuned RMSE", "RMSE Change"]].copy()
    add_table(slide, tuning_table, 0.75, 1.45, 5.9, 1.75, 10)
    add_body(
        slide,
        [
            "RandomizedSearchCV was applied to Random Forest and Gradient Boosting.",
            "The saved results show tuning did not improve chronological holdout RMSE.",
            "This is still useful: CV-optimized parameters do not always improve future-period performance.",
            "The baseline Random Forest remains the final recommended model for chronological prediction.",
        ],
        7.0,
        1.45,
        5.3,
        4.2,
        17,
    )

    slide = new_slide(prs, 14, "Decade-Wise Error")
    decade_table = error_by_decade[["Decade", "Rows", "MeanAbsoluteError", "RootMeanSquaredError"]].copy()
    add_table(slide, decade_table, 0.75, 1.45, 6.2, 2.4, 10)
    add_body(
        slide,
        [
            "RMSE increases from the 1990s to the 2010s.",
            "This suggests recent climate conditions are harder to predict using historical city-month patterns.",
            "The decade analysis adds interpretation beyond a single aggregate accuracy number.",
        ],
        7.3,
        1.65,
        4.9,
        3.2,
        17,
    )
    add_picture(slide, "anomaly_trend_by_decade.png", 1.0, 4.45, 11.1, 2.05)

    slide = new_slide(prs, 15, "South Asia Case Study")
    warming_table = warming[["Country", "WarmingRatePerDecade"]].head(6).copy()
    add_table(slide, warming_table, 0.75, 1.45, 4.7, 2.6, 10)
    add_picture(slide, "country_warming_rate_comparison.png", 5.75, 1.4, 3.25, 2.15)
    add_picture(slide, "pakistan_vs_global_trend.png", 9.25, 1.4, 3.0, 2.15)
    add_picture(slide, "pakistan_city_trends.png", 5.75, 4.1, 6.5, 2.2)
    add_body(
        slide,
        [
            "South Asia analysis gives the project regional climate relevance.",
            "Afghanistan and Pakistan have the highest estimated warming rates in the saved table.",
        ],
        0.85,
        4.45,
        4.45,
        1.6,
        15,
    )

    slide = new_slide(prs, 16, "Novelty and Contribution")
    add_body(
        slide,
        [
            "Novelty 1: anomaly-based climate interpretation, not only raw prediction.",
            "Novelty 2: chronological validation to reduce temporal leakage.",
            "Novelty 3: South Asia and Pakistan-focused analysis for local relevance.",
            "Novelty 4: comparative evaluation across baseline, tree, ensemble, and boosting models.",
            "Contribution: a reproducible city-level climate ML pipeline with real metrics, figures, and model artifacts.",
        ],
        0.75,
        1.55,
        6.4,
        4.9,
        18,
    )
    add_picture(slide, "rolling_mean_selected_cities.png", 7.25, 1.55, 5.1, 2.55)
    add_picture(slide, "global_vs_south_asia_trend.png", 7.25, 4.35, 5.1, 2.0)

    slide = new_slide(prs, 17, "Key Takeaways for Viva")
    add_body(
        slide,
        [
            "Why regression? AverageTemperature is continuous, so RMSE and R2 are appropriate.",
            "Why chronological split? Climate prediction should test future periods using earlier observations.",
            "Why MAPE is high? Near-zero temperatures inflate percentage error.",
            "What features mattered most? Historical city-month mean, lagged temperatures, and rolling means.",
            "What is novel? Anomaly analysis, time-aware validation, South Asia case study, and model comparison.",
        ],
        0.8,
        1.45,
        6.3,
        5.25,
        18,
    )
    add_metric(slide, "Best Model", str(best["Model"]), 7.6, 1.6, 3.2)
    add_metric(slide, "Chronological RMSE", f"{best['RMSE']:.3f}", 7.6, 3.05, 3.2)
    add_metric(slide, "Best CV Model", str(cv_summary.iloc[0]["Model"]), 7.6, 4.5, 3.2)

    slide = new_slide(prs, 18, "Conclusion and Future Work")
    add_body(
        slide,
        [
            f"Best chronological model: {best['Model']} with RMSE {best['RMSE']:.3f} and R2 {best['R2 Score']:.3f}.",
            f"Best CV model: {cv_summary.iloc[0]['Model']} with CV RMSE {cv_summary.iloc[0]['CV Mean RMSE']:.4f}.",
            "Lagged and historical city-month features are the strongest predictors.",
            "Chronological split gives a more realistic estimate than random split for climate forecasting.",
            "Future work: LSTM/Transformer models, post-2013 NASA/NOAA data, added CO2/rainfall/humidity variables, and a Streamlit app.",
        ],
        0.75,
        1.45,
        7.0,
        5.2,
        18,
    )
    add_metric(slide, "Final Model", str(best["Model"]), 8.4, 1.8, 3.1)
    add_metric(slide, "RMSE", f"{best['RMSE']:.3f}", 8.4, 3.25, 3.1)
    add_metric(slide, "R2", f"{best['R2 Score']:.3f}", 8.4, 4.7, 3.1)

    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
